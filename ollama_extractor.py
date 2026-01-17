# -*- coding: utf-8 -*-
"""
Windows app for image->text extraction and summarization using local Ollama.
- Mode A: Multi-file picker. Sort by ctime, extract per image with extraction model + prompt,
          then prepend a single summary using summary model + prompt. Save as Word-compatible HTML.
- Mode B (updated): Global hotkey (default: ctrl+alt+p) or "Snap Now" captures the ACTIVE WINDOW ONCE.
          Visual indication on snap. No background polling.
- Deferred processing (default ON): queue frames; only process + summarize on "Complete & Save…".
  If OFF: extract immediately per snap; summary still waits for "Complete & Save…".
- Stores prompts, selected models, hotkey, hotkey-enabled, defer toggle to config.json in the install directory.

Dependencies:
  pip install ollama PySide6 pillow keyboard pywin32 numpy

Notes:
  * Requires Ollama daemon to be running locally (http://localhost:11434). We do NOT start/stop it.
  * No OCR fallback; strictly VLM extraction as requested.
  * If you pick a non-vision model for extraction, it will error (as requested).
  * If context length is exceeded, request fails (as requested).
  * Windows 11 only target is fine.
"""

import os
import sys
import json
import base64
import re
import traceback
from pathlib import Path
from datetime import datetime

# UI
from PySide6.QtCore import Qt, QThread, Signal, QTimer
from PySide6.QtWidgets import (
    QApplication, QMainWindow, QWidget, QVBoxLayout, QHBoxLayout,
    QLabel, QPushButton, QFileDialog, QComboBox, QTextEdit, QSpinBox,
    QLineEdit, QMessageBox, QSystemTrayIcon, QMenu, QStyle, QCheckBox,
    QProgressBar, QDialog, QSlider, QGroupBox, QListWidget, QListWidgetItem,
    QFormLayout, QScrollArea, QFrame
)
from PySide6.QtGui import QAction, QPixmap, QPainter, QPen, QColor

# Imaging
from PIL import ImageGrab, Image
import numpy as np

# Global hotkey
import keyboard  # global hook (Windows); may require admin depending on system

# Active window capture (Windows)
import win32gui

# Ollama
import ollama
from ollama import Client, ResponseError

# OpenAI-compatible API (Azure, etc.)
from openai import OpenAI

# PowerPoint creation
from pptx import Presentation
from pptx.util import Inches


APP_NAME = "Visual Data Extractor"

# Word-friendly HTML output by default
DEFAULT_EXTRACTION_PROMPT = """You are producing a clean, Word-compatible HTML fragment for a single image.
Requirements:
- Do NOT include <html>, <head>, or <body>. Output only the inner HTML fragment.
- Structure the content using semantic HTML: <h2>, <h3>, <p>, <ul>, <li>, <table>, <thead>, <tbody>, <tr>, <th>, <td>, <code>, <pre>.
- Start with <h3>“Facts”</h3> followed by a <ul> of all key facts you can infer from the image.
- If there are numbers, dates, currencies, years, or measurements and you are confident you understand them include a <h3>“Data”</h3> section with a small <table> (two columns: Label, Value). 
- If applicable, Add <h3>“Narrative”</h3> with a short paragraph summarizing what’s in the image.

"""
DEFAULT_SUMMARY_PROMPT = """You are producing a clean, Word-compatible HTML fragment that is an executive summary of various data.
Requirements:
- Do NOT include <html>, <head>, or <body>. Output only the inner HTML fragment.
- Start with <h1>Executive Summary</h1>
- Provide a short <p> overview, then a <ul> of crisp bullet points (3–8 bullets).
- If there are recurring entities, dates, or numeric themes across images, include a compact <h2>Highlights</h2> section with a small <table> (three columns: Theme, Evidence, Notes).
- No markdown; HTML only. Keep it concise and executive-ready.
"""

DEFAULT_HOTKEY = "ctrl+alt+p"
CONFIG_FILE = Path(__file__).resolve().parent / "config.json"
IMAGE_EXTS = (".png", ".jpg", ".jpeg", ".webp", ".bmp", ".tif", ".tiff")
OLLAMA_HOST = "http://localhost:11434"  # local daemon only

# --- HTML helpers for Word-compatible output ---
HTML_CSS = """
<style>
  body { font-family: Segoe UI, Arial, sans-serif; line-height: 1.35; font-size: 11pt; color: #222; }
  h1, h2, h3 { margin: 0.4em 0 0.2em; }
  h1 { font-size: 20pt; }
  h2 { font-size: 16pt; border-bottom: 1px solid #ddd; padding-bottom: 4px; }
  h3 { font-size: 13pt; color: #333; }
  .meta { color: #666; font-size: 9pt; margin: 0.2em 0 0.8em; }
  ul { margin: 0.2em 0 0.8em 1.2em; }
  table { border-collapse: collapse; margin: 0.3em 0 0.8em; }
  th, td { border: 1px solid #ddd; padding: 6px 8px; }
  th { background: #f7f7f7; }
  pre { background: #fbfbfb; border: 1px solid #eee; padding: 8px; overflow-x: auto; }
  .section { margin: 1.0em 0; }
</style>
"""

def build_html_doc(title: str, summary_html: str, sections_html: list[str]) -> str:
    body = []
    if summary_html:
        body.append(summary_html)
    if sections_html:
        body.append("\n".join(sections_html))
    return f"""<!DOCTYPE html>
<html>
<head>
<meta charset="utf-8">
<title>{title}</title>
{HTML_CSS}
</head>
<body>
{''.join(body)}
</body>
</html>"""


def load_config():
    cfg = {
        "extraction_prompt": DEFAULT_EXTRACTION_PROMPT,
        "summary_prompt": DEFAULT_SUMMARY_PROMPT,
        "extraction_model": "",
        "summary_model": "",
        "hotkey": DEFAULT_HOTKEY,
        "hotkey_enabled": True,
        "defer_processing": True,
        "interim_save_folder": str(Path.home()),  # Default to user's home directory
        "include_summary": True,
        # Endpoint configuration
        "endpoints": [
            {"name": "Ollama (local)", "type": "ollama", "url": "http://localhost:11434", "api_key": ""}
        ],
        "selected_endpoint": "Ollama (local)",
    }
    if CONFIG_FILE.exists():
        try:
            cfg.update(json.loads(CONFIG_FILE.read_text(encoding="utf-8")))
        except Exception:
            # Ignore corrupt config, recreate later
            pass
    return cfg


def save_config(cfg: dict):
    try:
        CONFIG_FILE.write_text(json.dumps(cfg, indent=2), encoding="utf-8")
    except Exception:
        pass


def human_ts():
    return datetime.now().strftime("%Y-%m-%d %H:%M:%S")


def get_active_window_bbox():
    # Returns (left, top, right, bottom) for the foreground window
    hwnd = win32gui.GetForegroundWindow()
    if not hwnd:
        raise RuntimeError("No active window handle.")
    rect = win32gui.GetWindowRect(hwnd)
    return rect, hwnd


def detect_teams_crop_region(img: Image.Image) -> tuple[tuple[int, int, int, int], float] | None:
    """
    Detect Teams toolbar and side panel, return (crop_box, confidence).
    crop_box is (left, top, right, bottom) for PIL crop.
    confidence is 0.0-1.0 indicating detection certainty.
    Returns None if image is too small or detection completely fails.
    """
    width, height = img.size
    if width < 400 or height < 300:
        return None  # Image too small for meaningful detection

    # Convert to numpy array for analysis
    arr = np.array(img.convert('RGB'))

    # --- Detect top toolbar ---
    # Teams toolbar is typically dark (low brightness) and uniform
    # Scan from top to find where content begins

    # Convert to grayscale for brightness analysis
    gray = np.mean(arr, axis=2)

    # Analyze row-by-row brightness and variance
    top_crop = 0
    min_toolbar_height = 30
    max_toolbar_height = min(150, height // 4)

    for y in range(min_toolbar_height, max_toolbar_height):
        row_brightness = np.mean(gray[y])
        row_variance = np.var(gray[y])
        prev_brightness = np.mean(gray[y - 1])

        # Look for transition from dark/uniform to lighter/varied content
        # Teams toolbar is usually dark (brightness < 80) and uniform (variance < 500)
        if prev_brightness < 80 and row_variance < 500:
            if row_brightness > prev_brightness + 20 or row_variance > 800:
                top_crop = y
                break

    # --- Detect right side panel ---
    # The participant panel has high variance (many faces/colors)
    # The content area is typically more uniform (slides/documents)

    right_crop = width
    min_panel_width = width // 10  # At least 10% of width
    max_panel_width = width * 4 // 10  # At most 40% of width

    # Analyze vertical strips from right side
    strip_width = max(10, width // 50)
    variances = []

    for x in range(width - min_panel_width, width - max_panel_width, -strip_width):
        strip = arr[:, max(0, x - strip_width):x, :]
        strip_var = np.var(strip)
        variances.append((x, strip_var))

    if variances:
        # Find where variance changes significantly (content vs participant panel boundary)
        avg_variance = np.mean([v[1] for v in variances])

        # Look for edge detection - strong vertical edges indicate panel boundary
        # Compute horizontal gradient at each x position
        gray_arr = gray
        best_edge_x = None
        best_edge_score = 0

        for x in range(width - min_panel_width, width - max_panel_width, -strip_width):
            if x < 5 or x >= width - 5:
                continue
            # Compute vertical edge strength at this x
            left_col = gray_arr[:, x - 2:x]
            right_col = gray_arr[:, x:x + 2]
            edge_strength = np.mean(np.abs(np.mean(right_col, axis=1) - np.mean(left_col, axis=1)))

            # Also check if right side has higher variance than left
            left_region = arr[:, x - 50:x, :] if x > 50 else arr[:, :x, :]
            right_region = arr[:, x:x + 50, :] if x + 50 < width else arr[:, x:, :]
            left_var = np.var(left_region)
            right_var = np.var(right_region)

            # Score: strong edge + right side more varied than left
            if right_var > left_var * 1.2:  # Right side should be more varied
                score = edge_strength * (right_var / max(left_var, 1))
                if score > best_edge_score:
                    best_edge_score = score
                    best_edge_x = x

        if best_edge_x and best_edge_score > 5:
            right_crop = best_edge_x

    # Calculate confidence based on how clearly we detected boundaries
    confidence = 0.5  # Base confidence

    # Boost confidence if we found a clear toolbar
    if top_crop > min_toolbar_height:
        confidence += 0.2

    # Boost confidence if we found a clear side panel boundary
    if right_crop < width - min_panel_width:
        confidence += 0.3

    # Reduce confidence if crops are too aggressive
    content_width = right_crop
    content_height = height - top_crop
    if content_width < width * 0.5 or content_height < height * 0.7:
        confidence *= 0.5  # Suspicious - might be wrong

    # Build crop box (left, top, right, bottom)
    crop_box = (0, top_crop, right_crop, height)

    # Return None if we didn't detect anything useful
    if top_crop == 0 and right_crop == width:
        return None

    return (crop_box, min(confidence, 1.0))


def get_fallback_teams_crop(img: Image.Image, right_percent: int = 30, top_percent: int = 8) -> tuple[int, int, int, int]:
    """Get a fallback crop box removing top X% and right Y% of image."""
    width, height = img.size
    top_crop = height * top_percent // 100  # Remove top X%
    right_crop = width * (100 - right_percent) // 100  # Remove right Y%
    return (0, top_crop, right_crop, height)


class OllamaClientWrapper:
    def __init__(self, host=OLLAMA_HOST, timeout: float = 600.0):
        # Create client with explicit timeout (default 10 minutes for vision models)
        import httpx
        # Use longer connect timeout and even longer read timeout for vision inference
        self.client = Client(host=host, timeout=httpx.Timeout(timeout, connect=30.0))
        self.timeout = timeout

    def ensure_daemon(self):
        try:
            _ = ollama.list()  # health check
        except Exception as e:
            raise RuntimeError(
                "Cannot connect to the Ollama daemon at http://localhost:11434. "
                "Start Ollama first and try again."
            ) from e

    def list_models(self):
        resp = ollama.list()
        names = []
        for m in resp.get("models", []):
            n = m.get("name") or m.get("model")
            if n:
                names.append(n)
        return sorted(names, key=str.lower)

    def _resize_image_if_needed(self, image_bytes: bytes, max_dim: int = 2048) -> bytes:
        """Resize image if larger than max_dim to speed up vision model processing."""
        from PIL import Image
        from io import BytesIO

        img = Image.open(BytesIO(image_bytes))
        w, h = img.size

        if w <= max_dim and h <= max_dim:
            return image_bytes  # No resize needed

        # Calculate new size maintaining aspect ratio
        if w > h:
            new_w = max_dim
            new_h = int(h * (max_dim / w))
        else:
            new_h = max_dim
            new_w = int(w * (max_dim / h))

        img = img.resize((new_w, new_h), Image.LANCZOS)
        buf = BytesIO()
        img.save(buf, format="PNG", optimize=True)
        return buf.getvalue()

    def extract_from_image_path(self, model: str, user_prompt: str, image_path: str, temperature: float = 0.0) -> str:
        with open(image_path, "rb") as f:
            image_bytes = f.read()
        # Resize large images to speed up processing
        image_bytes = self._resize_image_if_needed(image_bytes)
        b64 = base64.b64encode(image_bytes).decode("utf-8")
        return self._extract_from_b64(model, user_prompt, b64, temperature)

    def extract_from_image_bytes(self, model: str, user_prompt: str, image_bytes: bytes, temperature: float = 0.0) -> str:
        # Resize large images to speed up processing
        image_bytes = self._resize_image_if_needed(image_bytes)
        b64 = base64.b64encode(image_bytes).decode("utf-8")
        return self._extract_from_b64(model, user_prompt, b64, temperature)

    def _extract_from_b64(self, model: str, user_prompt: str, b64: str, temperature: float) -> str:
        messages = [{
            "role": "user",
            "content": user_prompt,
            "images": [b64],
        }]
        resp = self.client.chat(model=model, messages=messages, options={"temperature": temperature}, stream=False)
        return resp["message"]["content"]

    def summarize_text(self, model: str, summary_prompt: str, text: str, temperature: float = 0.0) -> str:
        content = f"{summary_prompt}\n\n=== CONTENT START ===\n{text}\n=== CONTENT END ==="
        messages = [{"role": "user", "content": content}]
        resp = self.client.chat(model=model, messages=messages, options={"temperature": temperature}, stream=False)
        return resp["message"]["content"]


# Default Ollama endpoint
DEFAULT_OLLAMA_ENDPOINT = {
    "name": "Ollama (local)",
    "type": "ollama",
    "url": "http://localhost:11434",
    "api_key": ""
}


class AIClient:
    """Unified client for Ollama and OpenAI-compatible endpoints."""

    def __init__(self, endpoint_config: dict, timeout: float = 600.0):
        self.config = endpoint_config
        self.timeout = timeout
        self.endpoint_type = endpoint_config.get("type", "ollama")

        if self.endpoint_type == "ollama":
            self._init_ollama()
        else:
            self._init_openai()

    def _init_ollama(self):
        import httpx
        host = self.config.get("url", OLLAMA_HOST)
        self.ollama_client = Client(host=host, timeout=httpx.Timeout(self.timeout, connect=30.0))

    def _init_openai(self):
        self.openai_client = OpenAI(
            base_url=self.config.get("url", ""),
            api_key=self.config.get("api_key", ""),
            timeout=self.timeout
        )

    def check_connection(self) -> bool:
        """Check if the endpoint is reachable. Returns True if OK, raises exception otherwise."""
        try:
            if self.endpoint_type == "ollama":
                _ = ollama.list()
            else:
                # For OpenAI, try listing models
                _ = self.openai_client.models.list()
            return True
        except Exception as e:
            raise RuntimeError(f"Cannot connect to {self.config.get('name', 'endpoint')}: {e}") from e

    def list_models(self) -> list[str]:
        """Return list of available model names."""
        if self.endpoint_type == "ollama":
            resp = ollama.list()
            names = []
            for m in resp.get("models", []):
                n = m.get("name") or m.get("model")
                if n:
                    names.append(n)
            return sorted(names, key=str.lower)
        else:
            # OpenAI-compatible endpoint
            try:
                resp = self.openai_client.models.list()
                return sorted([m.id for m in resp.data], key=str.lower)
            except Exception:
                # If listing fails, return empty (user can type model name)
                return []

    def _resize_image_if_needed(self, image_bytes: bytes, max_dim: int = 2048) -> bytes:
        """Resize image if larger than max_dim to speed up vision model processing."""
        from io import BytesIO
        img = Image.open(BytesIO(image_bytes))
        w, h = img.size

        if w <= max_dim and h <= max_dim:
            return image_bytes

        if w > h:
            new_w = max_dim
            new_h = int(h * (max_dim / w))
        else:
            new_h = max_dim
            new_w = int(w * (max_dim / h))

        img = img.resize((new_w, new_h), Image.LANCZOS)
        buf = BytesIO()
        img.save(buf, format="PNG", optimize=True)
        return buf.getvalue()

    def extract_from_image_path(self, model: str, user_prompt: str, image_path: str, temperature: float = 0.0) -> str:
        with open(image_path, "rb") as f:
            image_bytes = f.read()
        return self.extract_from_image_bytes(model, user_prompt, image_bytes, temperature)

    def extract_from_image_bytes(self, model: str, user_prompt: str, image_bytes: bytes, temperature: float = 0.0) -> str:
        image_bytes = self._resize_image_if_needed(image_bytes)
        b64 = base64.b64encode(image_bytes).decode("utf-8")

        if self.endpoint_type == "ollama":
            return self._ollama_vision(model, user_prompt, b64, temperature)
        else:
            return self._openai_vision(model, user_prompt, b64, temperature)

    def _ollama_vision(self, model: str, user_prompt: str, b64: str, temperature: float) -> str:
        messages = [{
            "role": "user",
            "content": user_prompt,
            "images": [b64],
        }]
        resp = self.ollama_client.chat(model=model, messages=messages, options={"temperature": temperature}, stream=False)
        return resp["message"]["content"]

    def _openai_vision(self, model: str, user_prompt: str, b64: str, temperature: float) -> str:
        messages = [{
            "role": "user",
            "content": [
                {"type": "text", "text": user_prompt},
                {"type": "image_url", "image_url": {"url": f"data:image/png;base64,{b64}"}}
            ]
        }]
        resp = self.openai_client.chat.completions.create(
            model=model,
            messages=messages,
            temperature=temperature,
            max_tokens=4096
        )
        return resp.choices[0].message.content

    def summarize_text(self, model: str, summary_prompt: str, text: str, temperature: float = 0.0) -> str:
        content = f"{summary_prompt}\n\n=== CONTENT START ===\n{text}\n=== CONTENT END ==="

        if self.endpoint_type == "ollama":
            messages = [{"role": "user", "content": content}]
            resp = self.ollama_client.chat(model=model, messages=messages, options={"temperature": temperature}, stream=False)
            return resp["message"]["content"]
        else:
            messages = [{"role": "user", "content": content}]
            resp = self.openai_client.chat.completions.create(
                model=model,
                messages=messages,
                temperature=temperature,
                max_tokens=4096
            )
            return resp.choices[0].message.content


class ProcessingWorker(QThread):
    """Background worker for AI processing operations."""
    # Signals
    progress = Signal(str, int, int)  # (message, current, total)
    status = Signal(str)  # status message
    finished_ok = Signal(object)  # result data
    finished_error = Signal(str)  # error message

    def __init__(self, parent=None):
        super().__init__(parent)
        self._cancelled = False
        self.task = None
        self.task_data = None
        self.ai_client: AIClient | None = None

    def cancel(self):
        """Request cancellation of the current operation."""
        self._cancelled = True

    def is_cancelled(self):
        return self._cancelled

    def setup_extract_images(self, ai_client: AIClient, images: list, extract_model: str,
                              extraction_prompt: str, summary_model: str, summary_prompt: str,
                              include_summary: bool, image_type: str = "path"):
        """Setup task to extract from images (Mode A or queued frames)."""
        self.task = "extract_images"
        self.ai_client = ai_client
        self.task_data = {
            "images": images,  # list of paths or (timestamp, bytes) tuples
            "extract_model": extract_model,
            "extraction_prompt": extraction_prompt,
            "summary_model": summary_model,
            "summary_prompt": summary_prompt,
            "include_summary": include_summary,
            "image_type": image_type,  # "path" or "bytes"
        }
        self._cancelled = False

    def setup_extract_single(self, ai_client: AIClient, image_bytes: bytes,
                              extract_model: str, extraction_prompt: str, timestamp: str):
        """Setup task to extract from a single image (immediate Mode B)."""
        self.task = "extract_single"
        self.ai_client = ai_client
        self.task_data = {
            "image_bytes": image_bytes,
            "extract_model": extract_model,
            "extraction_prompt": extraction_prompt,
            "timestamp": timestamp,
        }
        self._cancelled = False

    def setup_rename_folder(self, ai_client: AIClient, summary_model: str,
                             content: str, folder_path: Path):
        """Setup task to generate folder name."""
        self.task = "rename_folder"
        self.ai_client = ai_client
        self.task_data = {
            "summary_model": summary_model,
            "content": content,
            "folder_path": folder_path,
        }
        self._cancelled = False

    def run(self):
        try:
            if self.task == "extract_images":
                self._run_extract_images()
            elif self.task == "extract_single":
                self._run_extract_single()
            elif self.task == "rename_folder":
                self._run_rename_folder()
        except Exception as e:
            self.finished_error.emit(f"{type(e).__name__}: {e}")

    def _run_extract_images(self):
        data = self.task_data
        images = data["images"]
        extract_model = data["extract_model"]
        extraction_prompt = data["extraction_prompt"]
        summary_model = data["summary_model"]
        summary_prompt = data["summary_prompt"]
        include_summary = data["include_summary"]
        image_type = data["image_type"]

        total = len(images) + (1 if include_summary else 0)
        sections = []
        all_extractions = []

        for i, img in enumerate(images, 1):
            if self._cancelled:
                self.finished_error.emit("Operation cancelled")
                return

            if image_type == "path":
                fname = Path(img).name
                self.progress.emit(f"Extracting image {i} of {len(images)}: {fname}", i, total)
                self.status.emit(f"Sending image to {extract_model}...")
                cts = datetime.fromtimestamp(os.path.getctime(img))
                fragment = self.ai_client.extract_from_image_path(extract_model, extraction_prompt, img).strip()
                section_html = f'''
<section class="section">
  <h2>{fname}</h2>
  <p class="meta">Created: {cts}</p>
  {fragment}
</section>
'''
            else:  # bytes - tuple of (timestamp, png_bytes) or (timestamp, png_bytes, path)
                when = img[0]
                png_bytes = img[1]
                self.progress.emit(f"Extracting queued frame {i} of {len(images)}", i, total)
                self.status.emit(f"Sending frame to {extract_model}...")
                fragment = self.ai_client.extract_from_image_bytes(extract_model, extraction_prompt, png_bytes).strip()
                section_html = f'''
<section class="section">
  <h2>Active window</h2>
  <p class="meta">Captured: {when}</p>
  {fragment}
</section>
'''
            sections.append(section_html)
            all_extractions.append(fragment)

        summary_html = ""
        if include_summary and not self._cancelled:
            self.progress.emit(f"Generating executive summary...", len(images) + 1, total)
            self.status.emit(f"Sending content to {summary_model} for summary...")
            summary_input = "\n\n".join(all_extractions)
            summary_html = self.ai_client.summarize_text(summary_model, summary_prompt, summary_input).strip()
            if not summary_html or len(summary_html) < 100:
                summary_html = '<div class="section"><h1>Executive Summary</h1><p>LLM did not return any summary data.</p></div>'
            elif not (summary_html.startswith('<') and ('<h1' in summary_html or '<section' in summary_html or '<div' in summary_html)):
                summary_html = f'<div class="section"><h1>Executive Summary</h1><p>{summary_html}</p></div>'

        if not self._cancelled:
            self.finished_ok.emit({
                "sections": sections,
                "all_extractions": all_extractions,
                "summary_html": summary_html,
            })

    def _run_extract_single(self):
        data = self.task_data
        self.status.emit(f"Sending image to {data['extract_model']}...")
        fragment = self.ai_client.extract_from_image_bytes(
            data["extract_model"], data["extraction_prompt"], data["image_bytes"]
        ).strip()

        section_html = f'''
<section class="section">
  <h2>Active window</h2>
  <p class="meta">Captured: {data["timestamp"]}</p>
  {fragment}
</section>
'''
        if not self._cancelled:
            self.finished_ok.emit({"section_html": section_html})

    def _run_rename_folder(self):
        data = self.task_data
        summary_model = data["summary_model"]
        content = data["content"]
        folder_path = data["folder_path"]

        if not folder_path or not folder_path.exists():
            self.finished_ok.emit({"suggested_name": None, "folder_path": None})
            return

        self.status.emit(f"Asking {summary_model} for a descriptive folder name...")

        naming_prompt = """Generate a very short folder name (2-4 words, lowercase, hyphens between words) that describes this content.
Rules:
- Use only lowercase letters, numbers, and hyphens
- No spaces, underscores, or special characters
- Maximum 30 characters
- Be descriptive but concise
- Output ONLY the folder name, nothing else

Example outputs: quarterly-sales-report, meeting-notes-jan, product-screenshots, bug-analysis"""

        raw_name = self.ai_client.summarize_text(summary_model, naming_prompt, content[:2000], temperature=0.3).strip()

        # Clean the name
        clean_name = raw_name.lower().strip()
        clean_name = re.sub(r'[^a-z0-9\-]', '-', clean_name)
        clean_name = re.sub(r'-+', '-', clean_name)
        clean_name = clean_name.strip('-')[:30]

        if not clean_name or len(clean_name) < 3:
            clean_name = "captures"

        date_suffix = datetime.now().strftime('%Y%m%d')
        suggested_name = f"{clean_name}_{date_suffix}"

        # Return suggested name without renaming - let the UI handle the rename after user confirmation
        self.finished_ok.emit({"suggested_name": suggested_name, "folder_path": folder_path})


class CropPreviewDialog(QDialog):
    """Dialog showing crop preview with adjustment controls and 3-button choice."""

    def __init__(self, pil_image: Image.Image, suggested_crop: tuple[int, int, int, int], parent=None, total_images: int = 1):
        super().__init__(parent)
        self.setWindowTitle("Adjust Crop Region")
        self.setModal(True)

        self.original_image = pil_image
        self.width, self.height = pil_image.size
        self.crop_box = list(suggested_crop)  # [left, top, right, bottom]
        self.total_images = total_images
        self._crop_mode = "none"  # "none", "all", or "each"

        # Scale factor for preview (max 800px wide)
        self.scale = min(1.0, 800 / self.width)
        self.preview_width = int(self.width * self.scale)
        self.preview_height = int(self.height * self.scale)

        self._setup_ui()
        self._update_preview()

    def _setup_ui(self):
        layout = QVBoxLayout(self)

        # Preview image label
        self.preview_label = QLabel()
        self.preview_label.setFixedSize(self.preview_width, self.preview_height)
        layout.addWidget(self.preview_label)

        # Crop controls group
        controls_group = QGroupBox("Crop Boundaries")
        controls_layout = QVBoxLayout(controls_group)

        # Top slider
        top_row = QHBoxLayout()
        top_row.addWidget(QLabel("Top:"))
        self.top_slider = QSlider(Qt.Horizontal)
        self.top_slider.setRange(0, self.height // 2)
        self.top_slider.setValue(self.crop_box[1])
        self.top_slider.valueChanged.connect(self._on_top_changed)
        top_row.addWidget(self.top_slider)
        self.top_value = QLabel(str(self.crop_box[1]))
        self.top_value.setMinimumWidth(50)
        top_row.addWidget(self.top_value)
        controls_layout.addLayout(top_row)

        # Right slider (inverted - higher value = more cropped)
        right_row = QHBoxLayout()
        right_row.addWidget(QLabel("Right edge:"))
        self.right_slider = QSlider(Qt.Horizontal)
        self.right_slider.setRange(self.width // 2, self.width)
        self.right_slider.setValue(self.crop_box[2])
        self.right_slider.valueChanged.connect(self._on_right_changed)
        right_row.addWidget(self.right_slider)
        self.right_value = QLabel(str(self.crop_box[2]))
        self.right_value.setMinimumWidth(50)
        right_row.addWidget(self.right_value)
        controls_layout.addLayout(right_row)

        layout.addWidget(controls_group)

        # Info label
        self.info_label = QLabel()
        self.info_label.setStyleSheet("color: #666;")
        layout.addWidget(self.info_label)
        self._update_info()

        # Buttons - 3-button choice
        btn_layout = QHBoxLayout()

        btn_no_crop = QPushButton("No Cropping")
        btn_no_crop.setToolTip("Use original image without any cropping")
        btn_no_crop.clicked.connect(self._on_no_cropping)

        # Button text depends on single vs multiple images
        if self.total_images <= 1:
            btn_apply_all = QPushButton("Apply Crop")
            btn_apply_all.setToolTip("Apply current crop settings to this image")
        else:
            btn_apply_all = QPushButton("Apply to All")
            btn_apply_all.setToolTip(f"Apply current crop settings to all {self.total_images} images")
        btn_apply_all.clicked.connect(self._on_apply_all)

        btn_adjust_each = QPushButton("Adjust Each Image")
        btn_adjust_each.setToolTip("Show this dialog for each image individually")
        btn_adjust_each.clicked.connect(self._on_adjust_each)

        # Only show "Adjust Each" if there are multiple images
        if self.total_images <= 1:
            btn_adjust_each.setVisible(False)

        btn_layout.addWidget(btn_no_crop)
        btn_layout.addWidget(btn_apply_all)
        btn_layout.addWidget(btn_adjust_each)
        layout.addLayout(btn_layout)

    def _on_top_changed(self, value):
        self.crop_box[1] = value
        self.top_value.setText(str(value))
        self._update_preview()
        self._update_info()

    def _on_right_changed(self, value):
        self.crop_box[2] = value
        self.right_value.setText(str(value))
        self._update_preview()
        self._update_info()

    def _update_info(self):
        orig_pixels = self.width * self.height
        crop_w = self.crop_box[2] - self.crop_box[0]
        crop_h = self.crop_box[3] - self.crop_box[1]
        crop_pixels = crop_w * crop_h
        pct = (crop_pixels / orig_pixels) * 100 if orig_pixels > 0 else 0
        self.info_label.setText(f"Cropped size: {crop_w} x {crop_h} ({pct:.0f}% of original)")

    def _update_preview(self):
        # Create a scaled preview with crop overlay
        preview = self.original_image.copy()
        preview = preview.resize((self.preview_width, self.preview_height), Image.LANCZOS)

        # Convert to QPixmap
        from io import BytesIO
        buf = BytesIO()
        preview.save(buf, format="PNG")
        buf.seek(0)
        pixmap = QPixmap()
        pixmap.loadFromData(buf.getvalue())

        # Draw crop rectangle overlay
        painter = QPainter(pixmap)

        # Draw darkened areas outside crop region
        pen = QPen(QColor(255, 0, 0, 200))
        pen.setWidth(2)
        painter.setPen(pen)

        # Scale crop coordinates to preview size
        scaled_left = int(self.crop_box[0] * self.scale)
        scaled_top = int(self.crop_box[1] * self.scale)
        scaled_right = int(self.crop_box[2] * self.scale)
        scaled_bottom = int(self.crop_box[3] * self.scale)

        # Draw semi-transparent overlay on cropped areas
        painter.setBrush(QColor(0, 0, 0, 100))
        painter.setPen(Qt.NoPen)
        # Top area
        painter.drawRect(0, 0, self.preview_width, scaled_top)
        # Right area
        painter.drawRect(scaled_right, scaled_top, self.preview_width - scaled_right, scaled_bottom - scaled_top)

        # Draw crop boundary
        painter.setPen(pen)
        painter.setBrush(Qt.NoBrush)
        painter.drawRect(scaled_left, scaled_top, scaled_right - scaled_left, scaled_bottom - scaled_top)

        painter.end()
        self.preview_label.setPixmap(pixmap)

    def _on_no_cropping(self):
        """User chose not to crop any images."""
        self._crop_mode = "none"
        self.accept()

    def _on_apply_all(self):
        """User wants to apply current crop to all images."""
        self._crop_mode = "all"
        self.accept()

    def _on_adjust_each(self):
        """User wants to adjust crop for each image individually."""
        self._crop_mode = "each"
        self.accept()

    def get_crop_region(self) -> tuple[int, int, int, int]:
        """Return the final crop region (left, top, right, bottom)."""
        return tuple(self.crop_box)

    def get_crop_mode(self) -> str:
        """Return the crop mode: 'none', 'all', or 'each'."""
        return self._crop_mode


class EndpointManagerDialog(QDialog):
    """Dialog for managing AI endpoints (Ollama, Azure, OpenAI-compatible)."""

    def __init__(self, endpoints: list[dict], parent=None):
        super().__init__(parent)
        self.setWindowTitle("Manage Endpoints")
        self.setModal(True)
        self.setMinimumWidth(500)

        # Copy endpoints list (we'll modify and return)
        self.endpoints = [ep.copy() for ep in endpoints]
        self._setup_ui()

    def _setup_ui(self):
        layout = QVBoxLayout(self)

        # Existing endpoints list
        layout.addWidget(QLabel("Saved Endpoints:"))
        self.endpoint_list = QListWidget()
        self.endpoint_list.setMaximumHeight(150)
        self._refresh_endpoint_list()
        layout.addWidget(self.endpoint_list)

        # Delete button
        delete_row = QHBoxLayout()
        delete_row.addStretch()
        self.btn_delete = QPushButton("Delete Selected")
        self.btn_delete.clicked.connect(self._delete_selected)
        delete_row.addWidget(self.btn_delete)
        layout.addLayout(delete_row)

        # Separator
        line = QFrame()
        line.setFrameShape(QFrame.HLine)
        line.setFrameShadow(QFrame.Sunken)
        layout.addWidget(line)

        # Add new endpoint section
        layout.addWidget(QLabel("Add New Endpoint:"))

        form = QFormLayout()
        self.name_edit = QLineEdit()
        self.name_edit.setPlaceholderText("e.g., Azure GPT-4o")
        form.addRow("Name:", self.name_edit)

        self.url_edit = QLineEdit()
        self.url_edit.setPlaceholderText("e.g., https://myazure.openai.azure.com/openai/v1/")
        form.addRow("URL:", self.url_edit)

        self.key_edit = QLineEdit()
        self.key_edit.setPlaceholderText("API key (will be stored locally)")
        self.key_edit.setEchoMode(QLineEdit.Password)
        form.addRow("API Key:", self.key_edit)

        layout.addLayout(form)

        add_row = QHBoxLayout()
        add_row.addStretch()
        self.btn_add = QPushButton("Add Endpoint")
        self.btn_add.clicked.connect(self._add_endpoint)
        add_row.addWidget(self.btn_add)
        layout.addLayout(add_row)

        # Close button
        layout.addSpacing(10)
        btn_row = QHBoxLayout()
        btn_row.addStretch()
        btn_close = QPushButton("Close")
        btn_close.clicked.connect(self.accept)
        btn_row.addWidget(btn_close)
        layout.addLayout(btn_row)

    def _refresh_endpoint_list(self):
        self.endpoint_list.clear()
        for ep in self.endpoints:
            item = QListWidgetItem(ep["name"])
            if ep["type"] == "ollama":
                item.setToolTip("Built-in Ollama endpoint (cannot be deleted)")
            else:
                item.setToolTip(f"URL: {ep['url']}")
            self.endpoint_list.addItem(item)

    def _delete_selected(self):
        current = self.endpoint_list.currentRow()
        if current < 0:
            QMessageBox.warning(self, "No selection", "Select an endpoint to delete.")
            return

        ep = self.endpoints[current]
        if ep["type"] == "ollama":
            QMessageBox.warning(self, "Cannot delete", "The built-in Ollama endpoint cannot be deleted.")
            return

        reply = QMessageBox.question(
            self, "Delete Endpoint",
            f"Delete endpoint '{ep['name']}'?",
            QMessageBox.Yes | QMessageBox.No
        )
        if reply == QMessageBox.Yes:
            del self.endpoints[current]
            self._refresh_endpoint_list()

    def _add_endpoint(self):
        name = self.name_edit.text().strip()
        url = self.url_edit.text().strip()
        key = self.key_edit.text().strip()

        if not name:
            QMessageBox.warning(self, "Missing name", "Please enter a name for the endpoint.")
            return
        if not url:
            QMessageBox.warning(self, "Missing URL", "Please enter the endpoint URL.")
            return
        if not key:
            QMessageBox.warning(self, "Missing API key", "Please enter the API key.")
            return

        # Check for duplicate name
        if any(ep["name"] == name for ep in self.endpoints):
            QMessageBox.warning(self, "Duplicate name", f"An endpoint named '{name}' already exists.")
            return

        # Add new endpoint
        self.endpoints.append({
            "name": name,
            "type": "openai",
            "url": url,
            "api_key": key
        })
        self._refresh_endpoint_list()

        # Clear fields
        self.name_edit.clear()
        self.url_edit.clear()
        self.key_edit.clear()

        QMessageBox.information(self, "Endpoint added", f"Endpoint '{name}' has been added.")

    def get_endpoints(self) -> list[dict]:
        """Return the updated list of endpoints."""
        return self.endpoints


class HotkeyListener(QThread):
    pressed = Signal()  # emitted on hotkey press

    def __init__(self, hotkey: str, parent=None):
        super().__init__(parent)
        self.hotkey = hotkey
        self._hook_id = None

    def run(self):
        try:
            self._hook_id = keyboard.add_hotkey(self.hotkey, lambda: self.pressed.emit())
            keyboard.wait()  # blocks in this thread
        except Exception:
            pass

    def rebind(self, hotkey: str):
        try:
            if self._hook_id is not None:
                keyboard.remove_hotkey(self._hook_id)
        except Exception:
            pass
        self.hotkey = hotkey
        try:
            self._hook_id = keyboard.add_hotkey(self.hotkey, lambda: self.pressed.emit())
        except Exception:
            pass

    def stop(self):
        try:
            if self._hook_id is not None:
                keyboard.remove_hotkey(self._hook_id)
        except Exception:
            pass


class MainWindow(QMainWindow):
    def __init__(self):
        super().__init__()
        self.setWindowTitle(APP_NAME)
        self.setMinimumWidth(760)

        self.cfg = load_config()

        # Initialize AI client with selected endpoint
        self.ai_client: AIClient | None = None
        self._init_ai_client()

        # ---------- State for Mode B ----------
        # If defer_processing=True: store raw PNG bytes to process upon "Complete".
        # If defer_processing=False: process immediately and store HTML sections here.
        self.deferred_queue: list[tuple[str, bytes]] = []       # (timestamp, png_bytes)
        self.processed_sections_html: list[str] = []            # HTML fragments already extracted
        self.current_session_folder: Path | None = None         # Session subfolder for current capture session

        # Background processing worker
        self.worker: ProcessingWorker | None = None
        self._pending_operation: str | None = None  # Track what operation is in progress
        self._pending_html: str | None = None  # Store HTML during rename operation
        self._pending_timestamp: str | None = None  # Store timestamp for immediate processing

        # ---------- UI ----------

        # === MAIN CONTROLS (always visible) ===

        # Hotkey controls
        self.hotkey_enabled_checkbox = QCheckBox("Hotkey enabled")
        self.hotkey_enabled_checkbox.setChecked(bool(self.cfg.get("hotkey_enabled", True)))
        self.hotkey_enabled_checkbox.setToolTip("Enable or disable the global hotkey for capturing")
        self.hotkey_enabled_checkbox.stateChanged.connect(self.on_hotkey_enabled_changed)

        self.hotkey_edit = QLineEdit(self.cfg.get("hotkey", DEFAULT_HOTKEY))
        self.hotkey_edit.setMaximumWidth(120)
        self.hotkey_edit.setToolTip("Global hotkey combination (e.g. ctrl+alt+p)")
        self.btn_rebind_hotkey = QPushButton("Bind")
        self.btn_rebind_hotkey.clicked.connect(self.rebind_hotkey)

        # Folder controls
        self.folder_label = QLabel(self._short_folder_path())
        self.folder_label.setStyleSheet("color: #555;")
        self.btn_set_folder = QPushButton("Change")
        self.btn_set_folder.clicked.connect(self.set_interim_save_folder)
        self.btn_open_folder = QPushButton("Open")
        self.btn_open_folder.clicked.connect(self.open_session_folder)

        # Main action buttons
        self.btn_create_ppt = QPushButton("Create PPT from Folder")
        self.btn_create_ppt.setToolTip("Select a folder of images, crop them, create PowerPoint")
        self.btn_create_ppt.clicked.connect(self.create_ppt_from_folder)

        self.btn_complete_ppt = QPushButton("Complete PPT")
        self.btn_complete_ppt.setVisible(False)
        self.btn_complete_ppt.setToolTip("Create PPT from captured frames")
        self.btn_complete_ppt.clicked.connect(self.complete_ppt_from_captures)

        # Status labels
        self.status_label = QLabel("Ready.")
        self.status_label.setWordWrap(True)
        self.queue_label = QLabel("Captured: 0 frames")
        self.queue_label.setStyleSheet("color:#555;")
        self.flash_label = QLabel("")
        self.flash_label.setStyleSheet("padding:4px 8px; color:#fff; background:#2e7d32; border-radius:6px;")
        self.flash_label.setVisible(False)

        # Progress bar
        self.activity_label = QLabel("")
        self.activity_label.setStyleSheet("font-weight: bold; color: #1565c0;")
        self.activity_label.setVisible(False)
        self.progress_bar = QProgressBar()
        self.progress_bar.setVisible(False)
        self.progress_bar.setTextVisible(True)
        self.progress_bar.setMinimum(0)
        self.progress_bar.setMaximum(100)
        self.btn_cancel = QPushButton("Cancel")
        self.btn_cancel.setVisible(False)
        self.btn_cancel.setStyleSheet("background-color: #c62828; color: white; font-weight: bold;")
        self.btn_cancel.clicked.connect(self.cancel_processing)

        # === AI ACTION BUTTONS (in main section) ===
        self.btn_mode_a = QPushButton("Process Saved Images")
        self.btn_mode_a.setToolTip("Select images and extract content using AI")
        self.btn_mode_a.clicked.connect(self.process_images_mode)

        self.btn_complete = QPushButton("Complete & Save HTML")
        self.btn_complete.setVisible(False)
        self.btn_complete.clicked.connect(self.complete_and_save)

        # === COLLAPSIBLE AI SETTINGS (toggle button + hidden container) ===
        self.btn_toggle_ai = QPushButton("▶ AI Settings")
        self.btn_toggle_ai.setStyleSheet("text-align: left; padding: 5px 10px; font-weight: bold;")
        self.btn_toggle_ai.setFlat(True)
        self.btn_toggle_ai.clicked.connect(self._toggle_ai_settings)

        # Container for AI settings (hidden by default)
        self.ai_settings_container = QWidget()
        self.ai_settings_container.setVisible(False)

        # Endpoint selection
        self.endpoint_combo = QComboBox()
        self._populate_endpoint_combo()
        self.endpoint_combo.currentTextChanged.connect(self.on_endpoint_changed)
        self.btn_manage_endpoints = QPushButton("Manage Endpoints...")
        self.btn_manage_endpoints.clicked.connect(self.manage_endpoints)

        # Model selection
        self.model_extract = QComboBox()
        self.model_extract.setToolTip("Vision model for image extraction")
        self.model_extract.setEditable(True)
        self.model_summary = QComboBox()
        self.model_summary.setToolTip("Model for summary generation")
        self.model_summary.setEditable(True)
        self.btn_refresh_models = QPushButton("Refresh")
        self.btn_refresh_models.clicked.connect(self.refresh_models)
        self.btn_test_connection = QPushButton("Test")
        self.btn_test_connection.setToolTip("Test connection using test_image.png")
        self.btn_test_connection.clicked.connect(self.test_ai_connection)

        # Prompts
        self.prompt_extract = QTextEdit()
        self.prompt_extract.setPlainText(self.cfg.get("extraction_prompt", DEFAULT_EXTRACTION_PROMPT))
        self.prompt_extract.setMaximumHeight(80)
        self.prompt_summary = QTextEdit()
        self.prompt_summary.setPlainText(self.cfg.get("summary_prompt", DEFAULT_SUMMARY_PROMPT))
        self.prompt_summary.setMaximumHeight(80)

        # AI options
        self.include_summary_checkbox = QCheckBox("Include Summary")
        self.include_summary_checkbox.setChecked(bool(self.cfg.get("include_summary", True)))
        self.include_summary_checkbox.stateChanged.connect(self.on_include_summary_changed)

        self.defer_checkbox = QCheckBox("Defer processing")
        self.defer_checkbox.setChecked(bool(self.cfg.get("defer_processing", True)))
        self.defer_checkbox.stateChanged.connect(self.on_defer_changed)

        # === LAYOUT ===
        root = QWidget()
        v = QVBoxLayout(root)

        # Hotkey row
        hotkey_row = QHBoxLayout()
        hotkey_row.addWidget(QLabel("Hotkey:"))
        hotkey_row.addWidget(self.hotkey_edit)
        hotkey_row.addWidget(self.btn_rebind_hotkey)
        hotkey_row.addWidget(self.hotkey_enabled_checkbox)
        hotkey_row.addStretch()
        v.addLayout(hotkey_row)

        # Folder row
        folder_row = QHBoxLayout()
        folder_row.addWidget(QLabel("Save folder:"))
        folder_row.addWidget(self.folder_label, 1)
        folder_row.addWidget(self.btn_set_folder)
        folder_row.addWidget(self.btn_open_folder)
        v.addLayout(folder_row)

        # Main actions row
        main_actions = QHBoxLayout()
        main_actions.addWidget(self.btn_create_ppt)
        main_actions.addWidget(self.btn_mode_a)
        main_actions.addWidget(self.btn_complete)
        main_actions.addStretch()
        main_actions.addWidget(self.queue_label)
        main_actions.addWidget(self.btn_complete_ppt)
        v.addLayout(main_actions)

        # Collapsible AI settings toggle
        v.addWidget(self.btn_toggle_ai)

        # AI settings container (hidden by default)
        ai_layout = QVBoxLayout(self.ai_settings_container)
        ai_layout.setContentsMargins(20, 5, 5, 10)  # Indent the settings

        # Endpoint row
        endpoint_row = QHBoxLayout()
        endpoint_row.addWidget(QLabel("Endpoint:"))
        endpoint_row.addWidget(self.endpoint_combo, 1)
        endpoint_row.addWidget(self.btn_manage_endpoints)
        ai_layout.addLayout(endpoint_row)

        # Models row
        models_row = QHBoxLayout()
        models_row.addWidget(QLabel("Extraction:"))
        models_row.addWidget(self.model_extract, 1)
        models_row.addWidget(QLabel("Summary:"))
        models_row.addWidget(self.model_summary, 1)
        models_row.addWidget(self.btn_refresh_models)
        models_row.addWidget(self.btn_test_connection)
        ai_layout.addLayout(models_row)

        # Prompts
        ai_layout.addWidget(QLabel("Extraction prompt:"))
        ai_layout.addWidget(self.prompt_extract)
        ai_layout.addWidget(QLabel("Summary prompt:"))
        ai_layout.addWidget(self.prompt_summary)

        # AI options row
        ai_options = QHBoxLayout()
        ai_options.addWidget(self.include_summary_checkbox)
        ai_options.addWidget(self.defer_checkbox)
        ai_options.addStretch()
        ai_layout.addLayout(ai_options)

        v.addWidget(self.ai_settings_container)

        # Progress section
        v.addWidget(self.activity_label)
        progress_row = QHBoxLayout()
        progress_row.addWidget(self.progress_bar, 1)
        progress_row.addWidget(self.btn_cancel)
        v.addLayout(progress_row)

        # Status
        v.addWidget(self.status_label)
        v.addWidget(self.flash_label, alignment=Qt.AlignLeft)

        v.addStretch()

        self.setCentralWidget(root)
        self.add_discard_button()

        # Tray icon
        self.tray = QSystemTrayIcon(self)
        icon = self.style().standardIcon(QStyle.SP_ComputerIcon)
        self.tray.setIcon(icon)
        self.tray.setToolTip(APP_NAME)
        tray_menu = QMenu()
        act_show = QAction("Show", self)
        act_show.triggered.connect(lambda: (self.showNormal(), self.raise_(), self.activateWindow()))
        act_snap = QAction("Snap Frame (active window)", self)
        act_snap.triggered.connect(self.on_hotkey_pressed)
        act_complete = QAction("Complete & Save…", self)
        act_complete.triggered.connect(self.complete_and_save)
        act_quit = QAction("Quit", self)
        act_quit.triggered.connect(QApplication.instance().quit)
        tray_menu.addAction(act_show)
        tray_menu.addSeparator()
        tray_menu.addAction(act_snap)
        tray_menu.addAction(act_complete)
        tray_menu.addSeparator()
        tray_menu.addAction(act_quit)
        self.tray.setContextMenu(tray_menu)
        self.tray.show()

        # Hotkey listener (optional, based on toggle)
        self.hotkey_listener: HotkeyListener | None = None
        if self.hotkey_enabled_checkbox.isChecked():
            self.start_hotkey_listener()

        # Populate models
        self.refresh_models()

        # Apply last used models if they exist
        ex = self.cfg.get("extraction_model", "")
        sm = self.cfg.get("summary_model", "")
        if ex:
            idx = self.model_extract.findText(ex)
            if idx >= 0: self.model_extract.setCurrentIndex(idx)
        if sm:
            idx = self.model_summary.findText(sm)
            if idx >= 0: self.model_summary.setCurrentIndex(idx)

        self.status_update(f"Ready. (Hotkey: {self.cfg.get('hotkey', DEFAULT_HOTKEY)})")
        self.update_queue_state()

    # --------------------------- Endpoint & AI Client helpers ---------------------------

    def _init_ai_client(self):
        """Initialize the AI client from the selected endpoint configuration."""
        selected_name = self.cfg.get("selected_endpoint", "Ollama (local)")
        endpoints = self.cfg.get("endpoints", [DEFAULT_OLLAMA_ENDPOINT])

        # Find the selected endpoint config
        endpoint_cfg = None
        for ep in endpoints:
            if ep.get("name") == selected_name:
                endpoint_cfg = ep
                break

        if endpoint_cfg is None and endpoints:
            endpoint_cfg = endpoints[0]

        if endpoint_cfg is None:
            endpoint_cfg = DEFAULT_OLLAMA_ENDPOINT

        try:
            self.ai_client = AIClient(endpoint_cfg)
        except Exception as e:
            QMessageBox.warning(self, "AI Client Error", f"Could not initialize AI client:\n{e}")
            self.ai_client = None

    def _short_folder_path(self) -> str:
        """Return a shortened folder path for display."""
        folder = self.cfg.get("interim_save_folder", str(Path.home()))
        path = Path(folder)
        if len(str(path)) > 50:
            return f".../{path.parent.name}/{path.name}"
        return str(path)

    def _populate_endpoint_combo(self):
        """Populate the endpoint dropdown from config."""
        # Block signals to prevent on_endpoint_changed from firing during populate
        self.endpoint_combo.blockSignals(True)
        try:
            self.endpoint_combo.clear()
            endpoints = self.cfg.get("endpoints", [DEFAULT_OLLAMA_ENDPOINT])
            for ep in endpoints:
                self.endpoint_combo.addItem(ep.get("name", "Unknown"))
            # Select the current endpoint
            selected = self.cfg.get("selected_endpoint", "Ollama (local)")
            idx = self.endpoint_combo.findText(selected)
            if idx >= 0:
                self.endpoint_combo.setCurrentIndex(idx)
        finally:
            self.endpoint_combo.blockSignals(False)

    def _toggle_ai_settings(self):
        """Toggle visibility of AI settings panel."""
        is_visible = self.ai_settings_container.isVisible()
        self.ai_settings_container.setVisible(not is_visible)
        if is_visible:
            self.btn_toggle_ai.setText("▶ AI Settings")
        else:
            self.btn_toggle_ai.setText("▼ AI Settings")

    def test_ai_connection(self):
        """Test the AI connection by extracting from test_image.png."""
        # Check if AI client is initialized
        if not self.ai_client:
            QMessageBox.warning(self, "No AI Client", "AI client is not initialized.")
            return

        # Check extraction model is selected
        extract_model = self.model_extract.currentText().strip()
        if not extract_model:
            QMessageBox.warning(self, "No Model", "Please select an extraction model first.")
            return

        # Find test image
        test_image_path = Path(__file__).resolve().parent / "test_image.png"
        if not test_image_path.exists():
            QMessageBox.warning(self, "Test Image Missing",
                f"test_image.png not found in:\n{test_image_path.parent}")
            return

        # Get extraction prompt
        extraction_prompt = self.prompt_extract.toPlainText().strip() or DEFAULT_EXTRACTION_PROMPT

        # Update UI
        self.btn_test_connection.setEnabled(False)
        self.btn_test_connection.setText("Testing...")
        self.status_update(f"Testing connection with {extract_model}...")
        QApplication.processEvents()

        try:
            # Attempt extraction
            result = self.ai_client.extract_from_image_path(
                extract_model, extraction_prompt, str(test_image_path)
            )

            # Show success dialog with result
            self.btn_test_connection.setText("Test")
            self.btn_test_connection.setEnabled(True)
            self.status_update("Test successful!")

            # Create a scrollable dialog to show the result
            dialog = QDialog(self)
            dialog.setWindowTitle("Test Successful")
            dialog.setMinimumSize(600, 400)
            layout = QVBoxLayout(dialog)

            # Success header
            header = QLabel(f"<b>Endpoint:</b> {self.cfg.get('selected_endpoint', 'Unknown')}<br>"
                          f"<b>Model:</b> {extract_model}<br>"
                          f"<b>Status:</b> <span style='color:green'>SUCCESS</span>")
            header.setTextFormat(Qt.RichText)
            layout.addWidget(header)

            # Result text
            layout.addWidget(QLabel("<b>Extracted content:</b>"))
            result_text = QTextEdit()
            result_text.setPlainText(result.strip())
            result_text.setReadOnly(True)
            layout.addWidget(result_text)

            # Close button
            btn_close = QPushButton("Close")
            btn_close.clicked.connect(dialog.accept)
            layout.addWidget(btn_close)

            dialog.exec()

        except Exception as e:
            self.btn_test_connection.setText("Test")
            self.btn_test_connection.setEnabled(True)
            self.status_update("Test failed!")

            QMessageBox.critical(self, "Test Failed",
                f"<b>Endpoint:</b> {self.cfg.get('selected_endpoint', 'Unknown')}<br>"
                f"<b>Model:</b> {extract_model}<br>"
                f"<b>Status:</b> <span style='color:red'>FAILED</span><br><br>"
                f"<b>Error:</b><br>{str(e)}")

    def on_endpoint_changed(self, endpoint_name: str):
        """Handle endpoint selection change."""
        if not endpoint_name:
            return
        self.cfg["selected_endpoint"] = endpoint_name
        save_config(self.cfg)
        self._init_ai_client()
        # Refresh models for the new endpoint
        self.refresh_models()

    def manage_endpoints(self):
        """Open the endpoint management dialog."""
        endpoints = self.cfg.get("endpoints", [DEFAULT_OLLAMA_ENDPOINT])
        dialog = EndpointManagerDialog(endpoints, self)
        if dialog.exec() == QDialog.Accepted:
            # Update config with modified endpoints
            self.cfg["endpoints"] = dialog.get_endpoints()
            save_config(self.cfg)
            # Refresh the combo and reinitialize client
            self._populate_endpoint_combo()
            self._init_ai_client()
            self.refresh_models()

    def _process_images_with_crop_dialog(self, image_files: list[Path]) -> list[Image.Image] | None:
        """
        Show crop dialog for first image, then process all images with selected crop mode.
        Returns list of processed PIL images, or None if user cancels.
        """
        if not image_files:
            return None

        # Open first image to get crop settings
        first_img = Image.open(image_files[0])
        width, height = first_img.size

        # Show crop dialog for first image with 3-button choice
        suggested_crop = get_fallback_teams_crop(first_img, right_percent=30, top_percent=8)
        dialog = CropPreviewDialog(first_img, suggested_crop, self, total_images=len(image_files))
        dialog.setWindowTitle(f"Adjust Crop (1 of {len(image_files)} images)")

        result = dialog.exec()
        if result == QDialog.Rejected:
            return None

        crop_mode = dialog.get_crop_mode()  # "none", "all", or "each"
        crop_percentages = None

        if crop_mode == "all":
            crop_box = dialog.get_crop_region()
            crop_percentages = (crop_box[1] / height, crop_box[2] / width)
        elif crop_mode == "none":
            crop_percentages = (0.0, 1.0)

        # Process all images
        cropped_images = []
        for i, img_path in enumerate(image_files):
            self.status_update(f"Processing image {i+1}/{len(image_files)}: {img_path.name}")
            QApplication.processEvents()

            try:
                img = Image.open(img_path)
                w, h = img.size

                if crop_mode == "each":
                    # Show dialog for each image
                    suggested = get_fallback_teams_crop(img, right_percent=30, top_percent=8)
                    each_dialog = CropPreviewDialog(img, suggested, self, total_images=1)  # Single image mode
                    each_dialog.setWindowTitle(f"Adjust Crop ({i+1} of {len(image_files)} images)")
                    if each_dialog.exec() == QDialog.Accepted:
                        each_mode = each_dialog.get_crop_mode()
                        if each_mode != "none":
                            crop_box = each_dialog.get_crop_region()
                            img = img.crop(crop_box)
                elif crop_percentages and crop_percentages != (0.0, 1.0):
                    top_pct, right_pct = crop_percentages
                    crop_box = (0, int(h * top_pct), int(w * right_pct), h)
                    img = img.crop(crop_box)

                cropped_images.append(img)
            except Exception as e:
                QMessageBox.warning(self, "Image error", f"Could not process {img_path.name}:\n{e}")
                continue

        return cropped_images if cropped_images else None

    def complete_ppt_from_captures(self):
        """Create a PPT from captured frames in the current session."""
        if not self.current_session_folder or not self.current_session_folder.exists():
            QMessageBox.warning(self, "No captures", "No session folder with captures exists.")
            return

        # Find all capture images
        image_files = sorted(self.current_session_folder.glob("*_capture.png"))
        if not image_files:
            QMessageBox.warning(self, "No captures", f"No capture images found in:\n{self.current_session_folder}")
            return

        self.status_update(f"Processing {len(image_files)} captured frames for PPT...")

        # Process images with crop dialog
        cropped_images = self._process_images_with_crop_dialog(image_files)
        if cropped_images is None:
            self.status_update("PPT creation cancelled.")
            return

        # Create PPT
        self.status_update(f"Creating PowerPoint with {len(cropped_images)} slides...")
        QApplication.processEvents()

        prs = self._create_ppt_from_images(cropped_images)

        # Save dialog
        default_name = f"{self.current_session_folder.name}_slides.pptx"
        out_path, _ = QFileDialog.getSaveFileName(
            self,
            "Save PowerPoint",
            str(self.current_session_folder / default_name),
            "PowerPoint (*.pptx);;All Files (*.*)"
        )

        if out_path:
            prs.save(out_path)
            self.status_update(f"Saved: {out_path}")
            QMessageBox.information(self, "PPT Created", f"PowerPoint saved with {len(cropped_images)} slides.")
        else:
            self.status_update("Save canceled.")

    # --------------------------- Utility & UI helpers ---------------------------

    def status_update(self, msg: str):
        self.status_label.setText(msg)
        self.tray.setToolTip(f"{APP_NAME}\n{msg}")

    def progress_start(self, activity: str, total_steps: int = 0, cancellable: bool = True):
        """Show progress bar and activity label. If total_steps > 0, use determinate mode."""
        self.activity_label.setText(activity)
        self.activity_label.setVisible(True)
        self.progress_bar.setVisible(True)
        self.btn_cancel.setVisible(cancellable)
        # Disable action buttons while processing
        self.btn_mode_a.setEnabled(False)
        self.btn_complete.setEnabled(False)
        if total_steps > 0:
            self.progress_bar.setMaximum(total_steps)
            self.progress_bar.setValue(0)
            self.progress_bar.setFormat("%v / %m")
        else:
            # Indeterminate mode
            self.progress_bar.setMaximum(0)
            self.progress_bar.setValue(0)
        QApplication.processEvents()

    def progress_update(self, activity: str, current_step: int = 0):
        """Update activity text and progress value."""
        self.activity_label.setText(activity)
        if self.progress_bar.maximum() > 0:
            self.progress_bar.setValue(current_step)
        QApplication.processEvents()

    def progress_finish(self):
        """Hide progress bar and activity label."""
        self.activity_label.setVisible(False)
        self.progress_bar.setVisible(False)
        self.btn_cancel.setVisible(False)
        self.progress_bar.setValue(0)
        self.progress_bar.setMaximum(100)
        # Re-enable action buttons
        self.btn_mode_a.setEnabled(True)
        self.btn_complete.setEnabled(True)
        self._pending_operation = None
        QApplication.processEvents()

    def cancel_processing(self):
        """Cancel the current processing operation."""
        if self.worker and self.worker.isRunning():
            self.worker.cancel()
            self.status_update("Cancelling operation...")
            self.btn_cancel.setEnabled(False)
            self.btn_cancel.setText("Cancelling...")

    def get_or_create_session_folder(self) -> Path:
        """Get the current session folder, creating a new one if needed."""
        if self.current_session_folder is None:
            base_folder = Path(self.cfg.get("interim_save_folder", str(Path.home())))
            session_name = f"session_{datetime.now().strftime('%Y%m%d_%H%M%S')}"
            self.current_session_folder = base_folder / session_name
            self.current_session_folder.mkdir(parents=True, exist_ok=True)
            self.status_update(f"New session folder: {self.current_session_folder}")
        return self.current_session_folder

    def open_session_folder(self):
        """Open the current session folder in Windows Explorer."""
        if self.current_session_folder and self.current_session_folder.exists():
            os.startfile(self.current_session_folder)
        else:
            # Open the base interim folder if no session yet
            base_folder = Path(self.cfg.get("interim_save_folder", str(Path.home())))
            if base_folder.exists():
                os.startfile(base_folder)
            else:
                QMessageBox.information(self, "No folder", "No session folder exists yet. Capture some frames first.")

    def reset_session(self):
        """Reset the session folder (start fresh for next capture batch)."""
        self.current_session_folder = None

    def _apply_teams_crop(self, img: Image.Image) -> Image.Image | None:
        """
        Apply Teams auto-crop to the image.
        Returns cropped image, original image (if user skips), or None (if user cancels).
        """
        # Try automatic detection
        result = detect_teams_crop_region(img)

        if result is not None:
            crop_box, confidence = result
            # High confidence threshold for automatic cropping
            if confidence >= 0.7:
                # Auto-crop without asking
                return img.crop(crop_box)
            else:
                # Show preview dialog for confirmation/adjustment
                suggested_crop = crop_box
        else:
            # No detection - use fallback suggestion
            suggested_crop = get_fallback_teams_crop(img, right_percent=30, top_percent=8)

        # Show crop preview dialog
        dialog = CropPreviewDialog(img, suggested_crop, self)
        if dialog.exec() == QDialog.Accepted:
            if dialog.get_crop_mode() == "none":
                return img  # No cropping requested
            final_crop = dialog.get_crop_region()
            return img.crop(final_crop)
        else:
            # Dialog was closed/cancelled - return original
            return img

    def flash_ping(self, text="Captured frame"):
        self.flash_label.setText(text)
        self.flash_label.setVisible(True)
        self.tray.showMessage(APP_NAME, text, QSystemTrayIcon.Information, 1200)
        QTimer.singleShot(900, lambda: self.flash_label.setVisible(False))

    def save_current_config(self):
        self.cfg["extraction_prompt"] = self.prompt_extract.toPlainText().strip() or DEFAULT_EXTRACTION_PROMPT
        self.cfg["summary_prompt"] = self.prompt_summary.toPlainText().strip() or DEFAULT_SUMMARY_PROMPT
        self.cfg["hotkey"] = self.hotkey_edit.text().strip() or DEFAULT_HOTKEY
        self.cfg["hotkey_enabled"] = self.hotkey_enabled_checkbox.isChecked()
        self.cfg["defer_processing"] = self.defer_checkbox.isChecked()
        self.cfg["extraction_model"] = self.model_extract.currentText().strip()
        self.cfg["summary_model"] = self.model_summary.currentText().strip()
        self.cfg["include_summary"] = self.include_summary_checkbox.isChecked()
        save_config(self.cfg)

    def update_queue_state(self):
        if self.defer_checkbox.isChecked():
            n = len(self.deferred_queue)
        else:
            n = len(self.processed_sections_html)

        # Update capture count for main UI
        capture_count = len(self.deferred_queue) if self.current_session_folder else 0
        self.queue_label.setText(f"Captured: {capture_count} frames")
        self.btn_complete_ppt.setVisible(capture_count > 0)

        # Update AI section completion button
        if n <= 0:
            self.btn_complete.setVisible(False)
            self.btn_discard.setVisible(False)
        else:
            self.btn_complete.setVisible(True)
            self.btn_discard.setVisible(True)

    # --------------------------- Models ---------------------------

    def refresh_models(self):
        if not self.ai_client:
            QMessageBox.warning(self, "No AI Client", "AI client is not initialized. Check your endpoint configuration.")
            return
        try:
            models = self.ai_client.list_models()
        except Exception as e:
            endpoint_name = self.cfg.get("selected_endpoint", "Unknown")
            QMessageBox.critical(self, "Connection Error", f"Failed to list models from '{endpoint_name}'.\n\n{e}")
            return
        old_extract = self.model_extract.currentText()
        old_summary = self.model_summary.currentText()
        self.model_extract.clear()
        self.model_summary.clear()
        self.model_extract.addItems(models)
        self.model_summary.addItems(models)
        if old_extract:
            i = self.model_extract.findText(old_extract)
            if i >= 0: self.model_extract.setCurrentIndex(i)
        if old_summary:
            i = self.model_summary.findText(old_summary)
            if i >= 0: self.model_summary.setCurrentIndex(i)

    # --------------------------- Mode A ---------------------------

    def process_images_mode(self):
        self.save_current_config()
        files, _ = QFileDialog.getOpenFileNames(
            self,
            "Select image files",
            "",
            "Images (*.png *.jpg *.jpeg *.webp *.bmp *.tif *.tiff);;All Files (*.*)"
        )
        if not files:
            return

        files = [f for f in files if Path(f).suffix.lower() in IMAGE_EXTS]
        if not files:
            QMessageBox.warning(self, "No images", "No supported image files selected.")
            return

        files.sort(key=lambda p: os.path.getctime(p))

        extract_model = self.model_extract.currentText().strip()
        summary_model = self.model_summary.currentText().strip()
        if not extract_model or not summary_model:
            QMessageBox.warning(self, "Models required", "Pick both extraction and summary models.")
            return

        extraction_prompt = self.prompt_extract.toPlainText().strip() or DEFAULT_EXTRACTION_PROMPT
        summary_prompt = self.prompt_summary.toPlainText().strip() or DEFAULT_SUMMARY_PROMPT
        include_summary = self.include_summary_checkbox.isChecked()
        total_steps = len(files) + (1 if include_summary else 0)

        # Start background worker
        self._pending_operation = "mode_a"
        self.progress_start(f"Processing {len(files)} images...", total_steps)
        self.status_update(f"Processing {len(files)} images…")

        self.worker = ProcessingWorker(self)
        self.worker.setup_extract_images(
            self.ai_client, files, extract_model, extraction_prompt,
            summary_model, summary_prompt, include_summary, image_type="path"
        )
        self.worker.progress.connect(self._on_worker_progress)
        self.worker.status.connect(self._on_worker_status)
        self.worker.finished_ok.connect(self._on_mode_a_finished)
        self.worker.finished_error.connect(self._on_worker_error)
        self.worker.start()

    def _on_worker_progress(self, message: str, current: int, total: int):
        """Handle progress updates from worker."""
        self.progress_bar.setMaximum(total)
        self.progress_bar.setValue(current)
        self.activity_label.setText(message)

    def _on_worker_status(self, message: str):
        """Handle status updates from worker."""
        self.status_update(message)

    def _on_worker_error(self, error_msg: str):
        """Handle errors from worker."""
        self.progress_finish()
        self.btn_cancel.setText("Cancel")
        self.btn_cancel.setEnabled(True)
        if "cancelled" in error_msg.lower():
            self.status_update("Operation cancelled.")
        elif "timeout" in error_msg.lower() or "timed out" in error_msg.lower():
            QMessageBox.warning(
                self, "Timeout Error",
                f"The operation timed out. This can happen with:\n"
                f"- Large or complex images\n"
                f"- Slower vision models\n"
                f"- Heavy system load\n\n"
                f"Try:\n"
                f"- Using a smaller/faster model (e.g., qwen3-vl:4b)\n"
                f"- Processing fewer images at once\n"
                f"- Checking if Ollama is responsive\n\n"
                f"Error: {error_msg}"
            )
            self.status_update("Timeout - try a faster model or smaller images.")
        else:
            QMessageBox.critical(self, "Processing Error", error_msg)
            self.status_update("Ready.")

    def _on_mode_a_finished(self, result: dict):
        """Handle Mode A completion."""
        self.progress_finish()
        self.btn_cancel.setText("Cancel")
        self.btn_cancel.setEnabled(True)

        sections = result["sections"]
        summary_html = result["summary_html"]
        final_html = build_html_doc("Image Extractions", summary_html, sections)
        self.save_word_html_via_dialog(final_html)
        self.status_update("Ready.")

    def save_word_html_via_dialog(self, html_text: str):
        default_name = datetime.now().strftime("%Y%m%d_%H%M_run-summary.doc")
        # Default to session folder if it exists, otherwise use home directory
        if self.current_session_folder and self.current_session_folder.exists():
            default_path = str(self.current_session_folder / default_name)
        else:
            default_path = default_name
        out_path, _ = QFileDialog.getSaveFileName(
            self,
            "Save Word-compatible HTML",
            default_path,
            "Word-compatible HTML (*.doc *.html *.htm);;All Files (*.*)"
        )
        if out_path:
            Path(out_path).write_text(html_text, encoding="utf-8")
            self.status_update(f"Saved: {out_path}")
        else:
            self.status_update("Save canceled.")

    # --------------------------- PPT Mode ---------------------------

    def create_ppt_from_folder(self):
        """Create a PowerPoint from images in a folder, with optional auto-cropping."""
        # Pick folder (default to interim save folder)
        default_folder = self.cfg.get("interim_save_folder", str(Path.home()))
        folder = QFileDialog.getExistingDirectory(
            self,
            "Select folder with captured images",
            default_folder
        )
        if not folder:
            return

        folder_path = Path(folder)

        # Find all image files
        image_files = []
        for ext in IMAGE_EXTS:
            image_files.extend(folder_path.glob(f"*{ext}"))
            image_files.extend(folder_path.glob(f"*{ext.upper()}"))

        # Remove duplicates and sort by name (which includes timestamp)
        image_files = sorted(set(image_files), key=lambda p: p.name)

        if not image_files:
            QMessageBox.warning(self, "No images", f"No image files found in:\n{folder}")
            return

        self.status_update(f"Processing {len(image_files)} images for PPT...")

        # Process images with crop dialog
        cropped_images = self._process_images_with_crop_dialog(image_files)
        if cropped_images is None:
            self.status_update("PPT creation cancelled.")
            return

        # Create PPT
        self.status_update(f"Creating PowerPoint with {len(cropped_images)} slides...")
        QApplication.processEvents()

        prs = self._create_ppt_from_images(cropped_images)

        # Save dialog
        default_name = f"{folder_path.name}_slides.pptx"
        out_path, _ = QFileDialog.getSaveFileName(
            self,
            "Save PowerPoint",
            str(folder_path / default_name),
            "PowerPoint (*.pptx);;All Files (*.*)"
        )

        if out_path:
            prs.save(out_path)
            self.status_update(f"Saved: {out_path}")
            # Offer to open the folder
            if QMessageBox.question(
                self, "PPT Created",
                f"PowerPoint saved with {len(cropped_images)} slides.\n\nOpen containing folder?",
                QMessageBox.Yes | QMessageBox.No
            ) == QMessageBox.Yes:
                os.startfile(Path(out_path).parent)
        else:
            self.status_update("Save canceled.")

    def _create_ppt_from_images(self, images: list[Image.Image]) -> Presentation:
        """Create a PowerPoint presentation with one image per slide."""
        import tempfile

        prs = Presentation()
        # Set 16:9 widescreen dimensions
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # Use blank slide layout
        blank_layout = prs.slide_layouts[6]

        for img in images:
            slide = prs.slides.add_slide(blank_layout)

            # Save image to temp file (python-pptx needs file path or file-like object)
            with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
                img.save(tmp, format="PNG")
                tmp_path = tmp.name

            try:
                # Calculate dimensions to fill slide while maintaining aspect ratio
                img_width, img_height = img.size
                img_aspect = img_width / img_height
                slide_aspect = 13.333 / 7.5

                if img_aspect > slide_aspect:
                    # Image is wider - fit to width
                    width = Inches(13.333)
                    height = Inches(13.333 / img_aspect)
                else:
                    # Image is taller - fit to height
                    height = Inches(7.5)
                    width = Inches(7.5 * img_aspect)

                # Center the image on the slide
                left = (prs.slide_width - width) / 2
                top = (prs.slide_height - height) / 2

                slide.shapes.add_picture(tmp_path, left, top, width, height)
            finally:
                # Clean up temp file
                try:
                    os.unlink(tmp_path)
                except Exception:
                    pass

        return prs

    # --------------------------- Mode B (updated: snap once per hotkey) ---------------------------

    def on_hotkey_pressed(self):
        """Capture active window once. Depending on 'defer', either queue raw image or extract now."""
        try:
            (l, t, r, b), hwnd = get_active_window_bbox()
            bbox = (l, t, r, b)
            img = ImageGrab.grab(bbox=bbox)

            from io import BytesIO
            buf = BytesIO()
            img.save(buf, format="PNG")
            png_bytes = buf.getvalue()
            when = human_ts()
        except Exception as e:
            QMessageBox.critical(self, "Capture error", f"Failed to capture active window:\n{e}")
            return
        # Save interim PNG to session subfolder (for resilience/recovery)
        try:
            session_folder = self.get_or_create_session_folder()
            fname = session_folder / f"{datetime.now().strftime('%Y%m%d_%H%M%S')}_capture.png"
            Path(fname).write_bytes(png_bytes)
            self.status_update(f"Captured: {fname.name} (in {session_folder.name})")
        except Exception:
            # Non-fatal if write fails; continue
            fname = None
        self.flash_ping("Captured frame")

        extract_model = self.model_extract.currentText().strip()
        if not extract_model:
            QMessageBox.warning(self, "Extraction model needed", "Pick an extraction model first.")
            return

        extraction_prompt = self.prompt_extract.toPlainText().strip() or DEFAULT_EXTRACTION_PROMPT

        if self.defer_checkbox.isChecked():
            # Queue image for later processing
            # Store (when, png_bytes) but also include the on-disk path if available for tracing
            try:
                stored_path = str(fname) if fname else ""
            except Exception:
                stored_path = ""
            if stored_path:
                self.deferred_queue.append((when, png_bytes, stored_path))
            else:
                self.deferred_queue.append((when, png_bytes))
            self.status_update(f"Queued frame at {when} (deferred)")
            self.update_queue_state()
        else:
            # Process immediately in background; store the section HTML for later summary
            self._pending_operation = "immediate"
            self._pending_timestamp = when
            self.progress_start(f"Extracting captured frame with {extract_model}...")

            self.worker = ProcessingWorker(self)
            self.worker.setup_extract_single(self.ai_client, png_bytes, extract_model, extraction_prompt, when)
            self.worker.status.connect(self._on_worker_status)
            self.worker.finished_ok.connect(self._on_immediate_finished)
            self.worker.finished_error.connect(self._on_worker_error)
            self.worker.start()

    def _on_immediate_finished(self, result: dict):
        """Handle immediate extraction completion."""
        self.progress_finish()
        self.btn_cancel.setText("Cancel")
        self.btn_cancel.setEnabled(True)

        section_html = result["section_html"]
        self.processed_sections_html.append(section_html)
        self.status_update(f"Processed frame at {self._pending_timestamp} (immediate)")
        self._pending_timestamp = None
        self.update_queue_state()

    def complete_and_save(self):
        """Process any queued frames (if deferred), then summarize ALL processed sections and save."""
        summary_model = self.model_summary.currentText().strip()
        if not summary_model:
            QMessageBox.warning(self, "Summary model needed", "Pick a summary model to finalize.")
            return

        extract_model = self.model_extract.currentText().strip()
        if not extract_model:
            QMessageBox.warning(self, "Extraction model needed", "Pick an extraction model.")
            return

        summary_prompt = self.prompt_summary.toPlainText().strip() or DEFAULT_SUMMARY_PROMPT
        extraction_prompt = self.prompt_extract.toPlainText().strip() or DEFAULT_EXTRACTION_PROMPT
        include_summary = self.include_summary_checkbox.isChecked()

        # Check if we have anything to process
        has_queued = self.defer_checkbox.isChecked() and len(self.deferred_queue) > 0
        has_processed = len(self.processed_sections_html) > 0

        if not has_queued and not has_processed:
            QMessageBox.information(self, "Nothing to do", "No frames have been captured yet.")
            return

        # If we have queued frames, process them with the worker
        if has_queued:
            queued_count = len(self.deferred_queue)
            total_steps = queued_count + (1 if include_summary else 0)

            self._pending_operation = "mode_b"
            self.progress_start(f"Processing {queued_count} queued frames...", total_steps)
            self.status_update(f"Processing {queued_count} queued frames…")

            self.worker = ProcessingWorker(self)
            self.worker.setup_extract_images(
                self.ai_client, list(self.deferred_queue), extract_model, extraction_prompt,
                summary_model, summary_prompt, include_summary, image_type="bytes"
            )
            self.worker.progress.connect(self._on_worker_progress)
            self.worker.status.connect(self._on_worker_status)
            self.worker.finished_ok.connect(self._on_mode_b_finished)
            self.worker.finished_error.connect(self._on_worker_error)
            self.worker.start()
        else:
            # No queued frames, but we have processed sections - just need summary
            self._finalize_mode_b_with_existing_sections(include_summary, summary_model, summary_prompt)

    def _finalize_mode_b_with_existing_sections(self, include_summary: bool, summary_model: str, summary_prompt: str):
        """Finalize Mode B when we already have processed sections (no queued frames)."""
        if not self.processed_sections_html:
            QMessageBox.information(self, "Nothing to do", "No frames have been processed yet.")
            return

        if include_summary:
            # Need to generate summary in background
            self._pending_operation = "mode_b_summary"
            self.progress_start("Generating summary...", 1)

            joined = "\n".join(self.processed_sections_html)
            self.worker = ProcessingWorker(self)
            # Use a simplified approach: setup for summary generation only
            self.worker.setup_extract_images(
                self.ai_client, [], "", "",  # No images to extract
                summary_model, summary_prompt, True, image_type="path"
            )
            # Override with direct summary call by using the finished signal creatively
            # Actually, let's just do the summary inline since there's no image extraction
            self._on_mode_b_finished({
                "sections": self.processed_sections_html,
                "all_extractions": self.processed_sections_html,  # Use sections as extractions for summary
                "summary_html": "",  # Will be generated
            })
        else:
            # No summary needed, just save
            final_html = build_html_doc("Active Window Extractions", "", self.processed_sections_html)
            self._complete_mode_b_save(final_html, "\n".join(self.processed_sections_html))

    def _on_mode_b_finished(self, result: dict):
        """Handle Mode B completion."""
        self.progress_finish()
        self.btn_cancel.setText("Cancel")
        self.btn_cancel.setEnabled(True)

        sections = result["sections"]
        summary_html = result["summary_html"]

        # If we processed queued frames, add them to our existing processed sections
        if self._pending_operation == "mode_b":
            self.processed_sections_html.extend(sections)
            self.deferred_queue.clear()

        # Build final HTML
        all_sections = self.processed_sections_html if self._pending_operation == "mode_b" else sections
        joined = "\n".join(all_sections)
        final_html = build_html_doc("Active Window Extractions", summary_html, all_sections)

        # Rename folder and save
        self._complete_mode_b_save(final_html, joined)

    def _complete_mode_b_save(self, final_html: str, content_for_naming: str):
        """Complete Mode B by renaming folder and saving."""
        # Rename session folder (this will start another worker for the rename)
        self._start_rename_folder(content_for_naming, final_html)

    def _start_rename_folder(self, content: str, final_html: str):
        """Start folder rename operation."""
        summary_model = self.model_summary.currentText().strip()
        if not self.current_session_folder or not summary_model:
            # Skip rename, just save
            self.save_word_html_via_dialog(final_html)
            self._cleanup_mode_b()
            return

        self._pending_html = final_html
        self._pending_operation = "rename"
        self.progress_start("Generating folder name...", cancellable=False)

        self.worker = ProcessingWorker(self)
        self.worker.setup_rename_folder(self.ai_client, summary_model, content, self.current_session_folder)
        self.worker.status.connect(self._on_worker_status)
        self.worker.finished_ok.connect(self._on_rename_finished)
        self.worker.finished_error.connect(self._on_rename_error)
        self.worker.start()

    def _on_rename_finished(self, result: dict):
        """Handle folder rename completion - show dialog for user to edit name."""
        self.progress_finish()
        suggested_name = result.get("suggested_name")
        folder_path = result.get("folder_path")

        if suggested_name and folder_path and folder_path.exists():
            # Show dialog for user to edit the folder name
            from PySide6.QtWidgets import QInputDialog
            edited_name, ok = QInputDialog.getText(
                self,
                "Rename Session Folder",
                "Edit the folder name (or leave as-is):",
                text=suggested_name
            )

            if ok and edited_name.strip():
                # Clean user input
                clean_name = edited_name.strip()
                clean_name = re.sub(r'[<>:"/\\|?*]', '-', clean_name)  # Remove invalid path chars
                clean_name = clean_name.strip('-. ')[:50]  # Reasonable length limit

                if clean_name:
                    new_path = folder_path.parent / clean_name
                    # Handle collision
                    counter = 1
                    base_name = clean_name
                    while new_path.exists() and new_path != folder_path:
                        new_path = folder_path.parent / f"{base_name}_{counter}"
                        counter += 1

                    if new_path != folder_path:
                        try:
                            folder_path.rename(new_path)
                            self.current_session_folder = new_path
                            self.status_update(f"Session folder renamed to: {new_path.name}")
                        except Exception as e:
                            self.status_update(f"Could not rename folder: {e}")
                    else:
                        self.status_update("Folder name unchanged.")
            else:
                self.status_update("Folder rename skipped.")

        # Now save the HTML
        self.save_word_html_via_dialog(self._pending_html)
        self._cleanup_mode_b()

    def _on_rename_error(self, error_msg: str):
        """Handle folder rename error (non-fatal)."""
        self.progress_finish()
        self.status_update(f"Could not rename folder: {error_msg}")
        # Still save the HTML
        self.save_word_html_via_dialog(self._pending_html)
        self._cleanup_mode_b()

    def _cleanup_mode_b(self):
        """Clean up state after Mode B completion."""
        self._pending_html = None
        self.processed_sections_html.clear()
        self.deferred_queue.clear()
        self.reset_session()
        self.update_queue_state()
        self.status_update("Ready.")

    # --------------------------- Toggles & bindings ---------------------------

    def rebind_hotkey(self):
        from PySide6.QtWidgets import QInputDialog
        hk, ok = QInputDialog.getText(self, "Bind Hotkey", "Enter new hotkey (e.g. ctrl+alt+p) or press Esc to cancel:", text=self.hotkey_edit.text())
        if ok and hk.strip():
            self.hotkey_edit.setText(hk.strip())
            self.hotkey_listener.rebind(hk.strip())
            self.cfg["hotkey"] = hk.strip()
            save_config(self.cfg)
            self.status_update(f"Global hotkey bound to: {hk.strip()}")
        else:
            self.status_update("Hotkey binding canceled.")

    def start_hotkey_listener(self):
        try:
            if self.hotkey_listener:
                self.hotkey_listener.stop()
            self.hotkey_listener = HotkeyListener(self.cfg.get("hotkey", DEFAULT_HOTKEY))
            self.hotkey_listener.pressed.connect(self.on_hotkey_pressed)
            self.hotkey_listener.start()
            self.status_update(f"Hotkey enabled ({self.cfg.get('hotkey', DEFAULT_HOTKEY)})")
        except Exception:
            # non-fatal; user can still click "Snap Now"
            self.status_update("Hotkey could not be enabled (permissions?).")

    def stop_hotkey_listener(self):
        try:
            if self.hotkey_listener:
                self.hotkey_listener.stop()
                self.hotkey_listener = None
        except Exception:
            pass
        self.status_update("Hotkey disabled")

    def on_hotkey_enabled_changed(self, state):
        enabled = state == Qt.Checked
        if enabled:
            self.start_hotkey_listener()
        else:
            self.stop_hotkey_listener()
        self.save_current_config()

    def on_defer_changed(self, state):
        # Just persist. Queue semantics handled at snap/complete time.
        self.save_current_config()
        self.update_queue_state()

    def on_include_summary_changed(self, state):
        self.cfg["include_summary"] = bool(state == Qt.Checked)
        save_config(self.cfg)

    # --------------------------- Lifecycle ---------------------------

    def closeEvent(self, event):
        # Warn if processing is in progress
        if self.worker and self.worker.isRunning():
            reply = QMessageBox.question(
                self, "Processing in Progress",
                "A processing operation is still running. Cancel it and close?",
                QMessageBox.Yes | QMessageBox.No, QMessageBox.No
            )
            if reply == QMessageBox.No:
                event.ignore()
                return
            # Cancel the worker
            self.worker.cancel()
            self.worker.wait(2000)  # Wait up to 2 seconds

        self.save_current_config()
        try:
            self.stop_hotkey_listener()
        except Exception:
            pass
        return super().closeEvent(event)

    def save_extraction_prompt(self):
        """Persist the current extraction prompt to config.json and update status."""
        self.cfg["extraction_prompt"] = self.prompt_extract.toPlainText().strip() or DEFAULT_EXTRACTION_PROMPT
        save_config(self.cfg)
        self.status_update("Extraction prompt saved.")

    def save_summary_prompt(self):
        """Persist the current summary prompt to config.json and update status."""
        self.cfg["summary_prompt"] = self.prompt_summary.toPlainText().strip() or DEFAULT_SUMMARY_PROMPT
        save_config(self.cfg)
        self.status_update("Summary prompt saved.")

    def restore_default_extraction_prompt(self):
        self.prompt_extract.setPlainText(DEFAULT_EXTRACTION_PROMPT)
        self.cfg["extraction_prompt"] = DEFAULT_EXTRACTION_PROMPT
        save_config(self.cfg)
        self.status_update("Extraction prompt restored to default.")

    def restore_default_summary_prompt(self):
        self.prompt_summary.setPlainText(DEFAULT_SUMMARY_PROMPT)
        self.cfg["summary_prompt"] = DEFAULT_SUMMARY_PROMPT
        save_config(self.cfg)
        self.status_update("Summary prompt restored to default.")

    def add_discard_button(self):
        self.btn_discard = QPushButton("Discard Session")
        self.btn_discard.setVisible(False)
        self.btn_discard.setToolTip("Discard all captured frames in the current session and delete their PNG files.\n"
                                    "Previous session folders are preserved for recovery via 'Process Saved Images'.")
        self.btn_discard.clicked.connect(self.discard_selected_frame)
        # Add to layout after queue_label
        self.centralWidget().layout().insertWidget(self.centralWidget().layout().indexOf(self.queue_label) + 1, self.btn_discard)

    def discard_selected_frame(self):
        # Discard ALL frames: clear both queued raw frames and any processed sections.
        queued = len(self.deferred_queue)
        processed = len(self.processed_sections_html)
        self.deferred_queue.clear()
        self.processed_sections_html.clear()

        # Remove interim PNG files from current session folder only (best-effort)
        session_deleted = False
        try:
            if self.current_session_folder and self.current_session_folder.exists():
                for p in self.current_session_folder.glob("*_capture.png"):
                    try:
                        p.unlink()
                    except Exception:
                        pass
                # Try to remove the empty session folder
                try:
                    if not any(self.current_session_folder.iterdir()):
                        self.current_session_folder.rmdir()
                        session_deleted = True
                except Exception:
                    pass
        except Exception:
            pass

        # Reset session
        self.reset_session()

        # Force the UI into the empty state immediately
        try:
            self.queue_label.setText("No outstanding frames to process")
            self.btn_complete.setVisible(False)
            self.btn_discard.setVisible(False)
            QApplication.processEvents()
        except Exception:
            pass

        msg = f"Discarded frames: queued={queued}, processed={processed}."
        if session_deleted:
            msg += " Session folder removed."
        self.status_update(msg)

    def set_interim_save_folder(self):
        """Set the folder where interim HTML files will be saved."""
        folder = QFileDialog.getExistingDirectory(self, "Select Save Folder", self.cfg.get("interim_save_folder", ""))
        if folder:
            self.cfg["interim_save_folder"] = folder
            save_config(self.cfg)
            self.folder_label.setText(self._short_folder_path())
            self.status_update(f"Save folder set: {folder}")
        else:
            self.status_update("Save folder not changed.")

def main():
    app = QApplication(sys.argv)
    win = MainWindow()
    win.show()
    sys.exit(app.exec())

if __name__ == "__main__":
    main()
